import chardet
import os
import re
import logging 
from fastapi import FastAPI, File, UploadFile, Depends, HTTPException, WebSocket, WebSocketDisconnect, Request, APIRouter, Form, Path, Query
import threading
from fastapi.responses import StreamingResponse
import json
import asyncio
from fastapi.middleware.cors import CORSMiddleware
import httpx
import tempfile
import subprocess
from fastapi.responses import FileResponse
import boto3
from botocore.exceptions import ClientError, BotoCoreError
from pydantic import BaseModel
import textwrap
from functools import lru_cache
import boto3
from fastapi import UploadFile, File, Depends, HTTPException, status, Body
from contextlib import closing
from tempfile import gettempdir
from fastapi.security import OAuth2AuthorizationCodeBearer
from jose import jwt, jwk
from jose.exceptions import JWTError
from jose.utils import base64url_decode
import time
import tempfile
from dotenv import load_dotenv
import ast
import esprima
import json
from botocore.config import Config
import uuid
import requests
from PyPDF2 import PdfReader, PdfWriter
import io
import pikepdf
import PyPDF2
import atexit
import fitz  # PyMuPDF
import io
from fastapi import HTTPException
from fastapi.responses import StreamingResponse
import logging
import json
from cachetools import TTLCache
from jose import jwt, JWTError
from fastapi import HTTPException, Security, Depends
from fastapi.security import HTTPAuthorizationCredentials, HTTPBearer
from yt_dlp import YoutubeDL
from datetime import datetime
from pptx import Presentation
from typing import List, Dict, Optional
from PIL import Image
import math
import base64

MAX_RETRIES = 1

# Load the .env file
load_dotenv()

S3_BUCKET_NAME = os.getenv('TEXTBOOK_S3_BUCKET')
if not S3_BUCKET_NAME:
    raise ValueError("TEXTBOOK_S3_BUCKET environment variable is not set")

router = APIRouter()

# Add AWS Bedrock client initialization
bedrock = boto3.client(
    service_name='bedrock-runtime',
    region_name=os.getenv('AWS_DEFAULT_REGION'),
    aws_access_key_id=os.getenv('AWS_ACCESS_KEY_ID'),
    aws_secret_access_key=os.getenv('AWS_SECRET_ACCESS_KEY')
)

bedrock_agent = boto3.client(
    service_name='bedrock-agent-runtime',
    region_name=os.getenv('AWS_DEFAULT_REGION'),
    aws_access_key_id=os.getenv('AWS_ACCESS_KEY_ID'),
    aws_secret_access_key=os.getenv('AWS_SECRET_ACCESS_KEY'),
    config=Config(region_name=os.getenv('AWS_DEFAULT_REGION'))
)

# Add AWS Translate client initialization
translate = boto3.client('translate')

# Add this with your other AWS client initializations
transcribe = boto3.client('transcribe',
    region_name='us-east-1'
)

# Initialize Polly client
polly_client = boto3.client('polly', region_name='us-east-1')

# Initialize S3 client
s3_client = boto3.client('s3',
    region_name=os.getenv('AWS_DEFAULT_REGION'),
    aws_access_key_id=os.getenv('AWS_ACCESS_KEY_ID'),
    aws_secret_access_key=os.getenv('AWS_SECRET_ACCESS_KEY')
)

# Initialize AWS Textract client
textract_client = boto3.client('textract')  # Add this line

# Initialize AWS Rekognition client
rekognition = boto3.client('rekognition')  # Add this line

# Initialize logging
logging.basicConfig(level=logging.DEBUG)

logging.getLogger('botocore').setLevel(logging.INFO)
logging.getLogger('botocore.parsers').setLevel(logging.WARNING)

app = FastAPI()

# Allow your frontend origin
origins = [
    "http://localhost:3000",  # React frontend origin
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

app.include_router(router)

# Configure the OAuth2 scheme
oauth2_scheme = OAuth2AuthorizationCodeBearer(
    authorizationUrl=os.getenv("COGNITO_AUTHORIZATION_URL"),
    tokenUrl=os.getenv("COGNITO_TOKEN_URL")
)

# Add this function to fetch and cache the JWKS
@lru_cache()
def get_jwks():
    jwks_url = os.getenv("COGNITO_JWKS_URL")
    response = requests.get(jwks_url)
    return response.json()["keys"]

# Cognito configuration
REGION = os.getenv('AWS_DEFAULT_REGION')
USER_POOL_ID = os.getenv('COGNITO_USER_POOL_ID')
APP_CLIENT_ID = os.getenv('COGNITO_APP_CLIENT_ID')
JWKS_URL = os.getenv('COGNITO_JWKS_URL')

# Fetch the JWKS from Cognito
jwks = requests.get(JWKS_URL).json()['keys']

security = HTTPBearer()

def decode_token(token: str):
    try:
        # Get the kid from the headers prior to verification
        headers = jwt.get_unverified_headers(token)
        kid = headers['kid']
        
        # Search for the kid in the downloaded public keys
        key_index = -1
        for i in range(len(jwks)):
            if kid == jwks[i]['kid']:
                key_index = i
                break
        if key_index == -1:
            raise HTTPException(status_code=401, detail='Public key not found in jwks.json')
        
        # Construct the public key
        public_key = jwk.construct(jwks[key_index])
        
        # Get the last two sections of the token,
        # message and signature (encoded in base64)
        message, encoded_signature = str(token).rsplit('.', 1)
        
        # Decode the signature
        decoded_signature = base64url_decode(encoded_signature.encode('utf-8'))
        
        # Verify the signature
        if not public_key.verify(message.encode("utf8"), decoded_signature):
            raise HTTPException(status_code=401, detail='Signature verification failed')
        
        # Since we passed the verification, we can now safely
        # use the unverified claims
        claims = jwt.get_unverified_claims(token)
        
        # Additionally we can verify the token expiration
        if time.time() > claims['exp']:
            raise HTTPException(status_code=401, detail='Token is expired')
        
        # And the Audience  (use claims['client_id'] if verifying an access token)
        if claims['aud'] != APP_CLIENT_ID:
            raise HTTPException(status_code=401, detail='Token was not issued for this audience')
        
        # Now we can use the claims
        return claims['sub']  # This is the user ID
    except JWTError:
        raise HTTPException(status_code=401, detail='Invalid token')

def get_current_user(credentials: HTTPAuthorizationCredentials = Security(security)):
    try:
        token = credentials.credentials
        return decode_token(token)
    except AttributeError:
        raise HTTPException(status_code=401, detail="Invalid authorization credentials")
    except Exception as e:
        raise HTTPException(status_code=401, detail=str(e))

learningCategories = {
    "Visual": [
        "I find it easier to understand new information when it is presented in diagrams, charts, or graphs.",
        "I prefer learning new concepts by observing demonstrations.",
        "I often visualize concepts or problems in my mind to help me solve them.",
        "I use colors, symbols, or drawings when taking notes to help me organize my thoughts.",
        "I remember information better when I see it written down or displayed on a screen",
    ],
    "Auditory": [
        "I learn better when I listen to explanations rather than read them.",
        "I prefer to learn by listening to audio lectures or podcasts.",
        "I remember information better when I hear it spoken aloud.",
        "I use mnemonics or chants to help me memorize information.",
        "I learn best by discussing concepts with others.",
    ],
    "ReadingWriting": [
        "I understand new ideas best when I write them down.",
        "I learn best by reading textbooks or articles.",
        "I use flashcards or mind maps to help me memorize information.",
        "I prefer to learn by reading and writing rather than listening or watching.",
        "I use diagrams or charts to help me understand and remember information.",
    ],
    "Kinesthetic": [
        "I enjoy working with physical models or doing hands-on activities to learn.",
        "I learn best by doing experiments or practical activities.",
        "I use role-playing or simulations to help me understand and apply new concepts.",
        "I prefer to learn by solving real-world problems or puzzles.",
        "I use physical models or manipulatives to help me visualize and understand information.",
    ]
}

@app.post("/save-learning-profile")
async def save_learning_profile(
    profile: dict = Body(...),
    current_user: str = Depends(get_current_user)
):
    try:
        S3_BUCKET_NAME = os.getenv('TEXTBOOK_S3_BUCKET')
        
        # Generate the text-based learning profile description
        learning_profile_description = generate_learning_profile_description(profile['answers'], learningCategories)
        
        # Create a dictionary with both the original answers and the generated description
        full_profile = {
            "answers": profile['answers'],
            "description": learning_profile_description
        }
        
        profile_key = f"learning_profiles/{current_user}.json"
        
        s3_client.put_object(
            Bucket=S3_BUCKET_NAME,
            Key=profile_key,
            Body=json.dumps(full_profile),
            ContentType='application/json'
        )
        
        return {"message": "Learning profile saved successfully", "description": learning_profile_description}
    except Exception as e:
        logging.error(f"Failed to save learning profile to S3: {str(e)}")
        raise HTTPException(status_code=500, detail="Failed to save learning profile")

def generate_learning_profile_description(answers: dict, questionnaire: dict) -> str:
    prompt = f"""Based on the following questionnaire answers, generate a paragraph-long textual description of the user's learning style. The answers are organized by learning category (Visual, Auditory, ReadingWriting, Kinesthetic) and represent the user's agreement level with each statement (1: Strongly Disagree, 5: Strongly Agree).

Questionnaire:
{json.dumps(questionnaire, indent=2)}

Questionnaire answers:
{json.dumps(answers, indent=2)}

Please provide a comprehensive description of the user's learning style, highlighting their strengths and preferences across different learning modalities. The description should be informative and tailored to the individual based on their responses."""

    try:
        body = json.dumps({
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 300,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": prompt
                        }
                    ]
                }
            ],
            "temperature": 0.7,
            "top_p": 0.9,
        })

        response = bedrock.invoke_model_with_response_stream(
            modelId="us.anthropic.claude-3-5-haiku-20241022-v1:0",
            body=body
        )

        full_response = ""
        for event in response['body']:
            chunk = json.loads(event['chunk']['bytes'])
            if chunk['type'] == 'content_block_delta':
                full_response += chunk['delta'].get('text', '')

        return full_response.strip()
    except Exception as e:
        logging.error(f"Error generating learning profile description: {str(e)}")
        return "Error generating learning profile description"

async def query_knowledge_base(query: str, top_k: int = 3):
    try:
        response = bedrock_agent.retrieve_and_generate(
            input={'text': query},
            retrieveAndGenerateConfiguration={
                'knowledgeBaseConfiguration': {
                    'knowledgeBaseId': os.getenv('KNOWLEDGE_BASE_ID'),
                    'modelArn': os.getenv('MODEL_ARN')
                },
                'type': 'KNOWLEDGE_BASE'
            }
        )
        
        generated_answer = response.get('output', {}).get('text', '')
        return generated_answer

    except ClientError as e:
        error_code = e.response['Error']['Code']
        error_message = e.response['Error']['Message']
        logging.error(f"AWS ClientError in query_knowledge_base: {error_code} - {error_message}")
        logging.error(f"Full error response: {e.response}")
        return f"Error querying knowledge base: {error_code} - {error_message}"
    except Exception as e:
        logging.error(f"Unexpected error in query_knowledge_base: {str(e)}")
        return f"Unexpected error querying knowledge base: {str(e)}"

# Add this function to upload files to S3
def upload_file_to_s3(file_name, bucket, object_name=None):
    if object_name is None:
        object_name = file_name
    try:
        s3_client.upload_file(file_name, bucket, object_name)
    except ClientError as e:
        logging.error(e)
        return False
    return True

async def create_textbook_collections(
    file_id: str,
    book_title: str,
    toc_structure: List[Dict],
    s3_key: str,
    current_user: str
) -> Dict:
    try:
        collections = {
            "book_id": file_id,
            "section_collections": [],
            "chapter_collections": []
        }

        # Create collections for each section
        for chapter in toc_structure:
            chapter_sections = chapter.get('sections', [])
            chapter_collection_id = str(uuid.uuid4())
            
            # Create collections for individual sections
            section_collection_ids = []
            for section in chapter_sections:
                section_id = str(uuid.uuid4())
                section_collection = {
                    "collection_id": section_id,
                    "name": f"{book_title} - {section['title']}",
                    "created_date": datetime.now().isoformat(),
                    "user_id": current_user,
                    "parent_chapter": chapter['number'],
                    "materials": {
                        "textbook_sections": [{
                            "section_id": section_id,
                            "title": section['title'],
                            "page": section['page'],
                            "s3_key": s3_key,
                            "added_date": datetime.now().isoformat()
                        }],
                        "transcriptions": [],
                        "presentations": [],
                        "notes": []
                    }
                }
                
                # Save section collection
                section_collection_key = f"collections/{current_user}/{section_id}.json"
                s3_client.put_object(
                    Bucket=S3_BUCKET_NAME,
                    Key=section_collection_key,
                    Body=json.dumps(section_collection),
                    ContentType='application/json'
                )
                
                section_collection_ids.append(section_id)
                collections["section_collections"].append(section_collection)

            # Create collection for the chapter
            chapter_collection = {
                "collection_id": chapter_collection_id,
                "name": f"{book_title} - {chapter['number']}: {chapter['title']}",
                "created_date": datetime.now().isoformat(),
                "user_id": current_user,
                "chapter_number": chapter['number'],
                "materials": {
                    "textbook_sections": [],
                    "transcriptions": [],
                    "presentations": [],
                    "notes": [],
                    "subcollections": section_collection_ids  # Reference to section collections
                }
            }
            
            # Add all sections to the chapter collection
            for section in chapter_sections:
                chapter_collection["materials"]["textbook_sections"].append({
                    "section_id": str(uuid.uuid4()),
                    "title": section['title'],
                    "page": section['page'],
                    "s3_key": s3_key,
                    "added_date": datetime.now().isoformat()
                })
            
            # Save chapter collection
            chapter_collection_key = f"collections/{current_user}/{chapter_collection_id}.json"
            s3_client.put_object(
                Bucket=S3_BUCKET_NAME,
                Key=chapter_collection_key,
                Body=json.dumps(chapter_collection),
                ContentType='application/json'
            )
            
            collections["chapter_collections"].append(chapter_collection)

        return collections

    except Exception as e:
        logging.error(f"Error creating textbook collections: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail="Failed to create textbook collections"
        )

@app.post("/upload-pdf")
async def upload_pdf(
    file: UploadFile = File(...),
    documentType: str = Form(...),
    tocPages: str = Form(None),
    current_user: str = Depends(get_current_user)
):
    logging.info(f"Received PDF upload request. Document type: {documentType}")
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="File format not supported. Please upload a PDF file.")
    
    # Generate a unique filename
    unique_filename = f"{uuid.uuid4()}_{file.filename}"
    s3_key = f"user-uploads/{current_user}/{unique_filename}"
    
    # Save the file temporarily
    temp_file_path = f"/tmp/{unique_filename}"
    with open(temp_file_path, "wb") as buffer:
        buffer.write(await file.read())
    
    try:
        # Upload to S3
        s3_client.upload_file(temp_file_path, S3_BUCKET_NAME, s3_key)
        logging.info(f"File uploaded to S3: {s3_key}")
        
        metadata = {
            "title": file.filename,
            "s3_key": s3_key,
            "user_id": current_user,
            "document_type": documentType,
            "table_of_contents": []  # Initialize with an empty list
        }
        
        collections_info = None
        if documentType == 'textbook' and tocPages:
            logging.info(f"Processing textbook TOC pages: {tocPages}")
            file_id = unique_filename.split('_')[0]
            start_page, end_page = map(int, tocPages.split('-'))
            
            # Await the TOC processing
            toc_structure = await process_toc_pages(
                S3_BUCKET_NAME, 
                s3_key, 
                start_page, 
                end_page, 
                file_id, 
                file.filename, 
                current_user, 
                unique_filename
            )
            
            metadata["table_of_contents"] = toc_structure
            
            # Create collections for chapters and sections
            collections_info = await create_textbook_collections(
                file_id=file_id,
                book_title=file.filename,
                toc_structure=toc_structure,
                s3_key=s3_key,
                current_user=current_user
            )
            logging.info(f"Created collections for textbook: {collections_info}")
        
        # Save metadata in S3
        metadata_key = f"metadata/{current_user}/{unique_filename}.json"
        s3_client.put_object(
            Bucket=S3_BUCKET_NAME,
            Key=metadata_key,
            Body=json.dumps(metadata),
            ContentType='application/json'
        )
        
        # Clean up the temporary file
        os.remove(temp_file_path)
        
        response = {"message": "PDF uploaded successfully", "s3_key": s3_key}
        if collections_info:
            response["collections"] = collections_info
        return response
        
    except Exception as e:
        logging.error(f"Failed to upload PDF to S3: {str(e)}")
        raise HTTPException(status_code=500, detail="Failed to upload PDF")

async def prepare_toc_images(pdf_path: str, start_page: int, end_page: int) -> List[bytes]:
    """Convert PDF pages to 1152x1152 images."""
    try:
        images = []
        with fitz.open(pdf_path) as pdf:
            for page_num in range(start_page - 1, end_page):
                page = pdf[page_num]
                # Get the page's pixmap
                pix = page.get_pixmap()
                
                # Convert to PIL Image
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                
                # Resize to 1152x1152 maintaining aspect ratio
                target_size = (1152, 1152)
                
                # Calculate dimensions maintaining aspect ratio
                aspect_ratio = img.width / img.height
                if aspect_ratio > 1:
                    new_width = 1152
                    new_height = int(1152 / aspect_ratio)
                else:
                    new_height = 1152
                    new_width = int(1152 * aspect_ratio)
                
                # Resize image
                img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
                
                # Create new white background image
                background = Image.new('RGB', target_size, (255, 255, 255))
                
                # Calculate position to paste resized image
                paste_x = (target_size[0] - new_width) // 2
                paste_y = (target_size[1] - new_height) // 2
                
                # Paste resized image onto white background
                background.paste(img, (paste_x, paste_y))
                
                # Convert to bytes
                img_byte_arr = io.BytesIO()
                background.save(img_byte_arr, format='JPEG')
                img_byte_arr = img_byte_arr.getvalue()
                
                images.append(img_byte_arr)
        
        return images
    
    except Exception as e:
        logging.error(f"Error preparing TOC images: {str(e)}")
        raise

async def process_toc_with_claude(images: List[bytes]) -> Dict:
    """Process TOC images with Claude 3 Sonnet vision model."""
    try:
        # Start with initial instruction
        messages = [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": """Analyze this table of contents and output ONLY a JSON object with this exact structure:
                        {
                            "chapters": [
                                {
                                    "number": "Chapter 1",
                                    "title": "To deliver you from the Preliminary Terrors",
                                    "page": 1,
                                    "sections": []
                                }
                            ]
                        }
                        
                        Important rules:
                        1. For each chapter:
                           - "number" should be "Chapter X" (or "Prologue" for the prologue)
                           - Include the full title after the chapter number
                           - "page" should be the integer page number from the right column
                        2. Convert Roman numerals (I, II, III, etc) to regular numbers in the chapter titles
                        3. Include the Prologue as a chapter with "number": "Prologue"
                        4. Make sure all page numbers are integers
                        5. Do not include any explanatory text in your response, only the JSON object
                        6. Keep the exact key names shown in the example ("number", "title", "page", "sections")
                        
                        Example chapter entry:
                        {
                            "number": "Chapter 1",
                            "title": "To deliver you from the Preliminary Terrors",
                            "page": 1,
                            "sections": []
                        }"""
                    }
                ]
            },
            {
                "role": "assistant",
                "content": "I understand. I'll analyze the table of contents and output only the JSON object with the specified structure, maintaining the exact format for chapter numbers, titles, and page numbers."
            }
        ]

        # Combine all images into a single user message
        image_content = []
        for i, img_bytes in enumerate(images, 1):
            image_content.extend([
                {
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": "image/jpeg",
                        "data": base64.b64encode(img_bytes).decode('utf-8')
                    }
                },
                {
                    "type": "text",
                    "text": f"This is page {i} of the table of contents."
                }
            ])
        
        # Add all images in a single user message
        messages.append({
            "role": "user",
            "content": image_content
        })

        # Make the API request
        body = json.dumps({
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 2048,
            "messages": messages,
            "temperature": 0.7,
            "top_p": 0.9
        })

        # Create a new bedrock client with proper configuration
        bedrock = boto3.client(
            'bedrock-runtime',
            config=Config(
                retries={'max_attempts': 3},
                connect_timeout=30,
                read_timeout=30,
                user_agent_extra=f"inference_profile={os.getenv('CLAUDE_INFERENCE_PROFILE_ARN')}"
            )
        )

        response = bedrock.invoke_model(
            modelId="us.anthropic.claude-3-5-sonnet-20241022-v2:0",
            body=body,
            contentType="application/json",
            accept="application/json"
        )

        # Parse the response
        response_body = json.loads(response['body'].read())
        
        # Extract the content from Claude's response
        assistant_message = response_body.get('content', [{}])[0].get('text', '{}')
        
        # Parse the JSON from the text response
        toc_data = json.loads(assistant_message)
        
        logging.debug(f"Parsed TOC data: {toc_data}")  # Add debug logging
        return toc_data

    except Exception as e:
        logging.error(f"Error processing TOC with Claude: {str(e)}")
        logging.error(f"Full error: {str(e)}")
        raise

# Update the process_toc_pages function to use Claude instead of Llama
async def process_toc_pages(bucket: str, document_key: str, start_page: int, end_page: int, file_id: str, book_title: str, current_user: str, unique_filename: str) -> List[Dict]:
    """Process TOC pages and convert to chapter/section structure."""
    logging.info(f"Starting TOC processing for document: {document_key}")
    logging.info(f"TOC pages to process: {start_page}-{end_page}")
    
    try:
        # Download the PDF from S3
        response = s3_client.get_object(Bucket=bucket, Key=document_key)
        pdf_content = response['Body'].read()

        # Create a temporary file to store the downloaded PDF
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_file:
            temp_file.write(pdf_content)
            temp_path = temp_file.name

        # Prepare images for Claude model
        toc_images = await prepare_toc_images(temp_path, start_page, end_page)
        
        # Process images with Claude model
        toc_data = await process_toc_with_claude(toc_images)
        
        # The response from Claude is already in the correct format
        chapters = []
        for chapter in toc_data['chapters']:
            chapter_info = {
                'number': chapter['number'],  # Already in correct format "Chapter X"
                'title': chapter['title'],
                'page': chapter['page'],
                'sections': chapter.get('sections', [])  # Use get() with default empty list
            }
            chapters.append(chapter_info)

        return chapters

    except Exception as e:
        logging.error(f"Error processing TOC pages: {str(e)}")
        logging.error(f"Error type: {type(e).__name__}")
        logging.error(f"Error args: {e.args}")
        return []

    finally:
        # Clean up the temporary file
        if 'temp_path' in locals():
            os.unlink(temp_path)

def extract_chapters_from_textract(textract_response, content_start_page):
    chapters = []
    current_chapter = None
    current_section = None
    
    for block in textract_response['Blocks']:
        if block['BlockType'] == 'LINE':
            text = block['Text'].strip()
            logging.debug(f"Processing line: {text}")
            
            if text.lower().startswith('chapter ') and re.match(r'^\d+(\.\d+)?:?', text.split()[1]):
                logging.info(f"Found new chapter: {text}")
                if current_chapter:
                    chapters.append(current_chapter)
                chapter_parts = text.split(':', 1)
                current_chapter = {
                    'number': chapter_parts[0].strip(),
                    'title': chapter_parts[1].strip() if len(chapter_parts) > 1 else '',
                    'sections': []
                }
            elif current_chapter and re.match(r'^\d+\.\d+', text):
                logging.info(f"Found potential section: {text}")
                current_section = {
                    'title': text,
                    'page': None
                }
            elif current_section and text.isdigit():
                # Adjust the page number by adding the content_start_page
                adjusted_page = int(text) + content_start_page - 1
                current_section['page'] = adjusted_page
                logging.debug(f"Created section: {current_section}")
                if current_section['title'] and current_section['page']:
                    current_chapter['sections'].append(current_section)
                    logging.info(f"Added section: {current_section} to chapter {current_chapter['number']}")
                current_section = None
    
    if current_chapter:
        chapters.append(current_chapter)
    
    # Detailed logging of the final structure
    logging.info("Final chapter structure:")
    for chapter in chapters:
        logging.info(f"Chapter {chapter['number']}: {chapter['title']}")
        logging.info(f"  Sections: {chapter['sections']}")
    
    total_sections = sum(len(ch['sections']) for ch in chapters)
    logging.info(f"Extracted {len(chapters)} chapters with a total of {total_sections} sections")
    return chapters

def get_text_from_block(textract_response, block):
    if 'Text' in block:
        return block['Text']
    elif 'Relationships' in block:
        child_blocks = [b for b in textract_response['Blocks'] if b['Id'] in block['Relationships'][0]['Ids']]
        return ' '.join([get_text_from_block(textract_response, child) for child in child_blocks])
    return ''

@app.get("/user-books")
async def get_user_books(current_user: str = Depends(get_current_user)):
    S3_BUCKET_NAME = os.getenv('TEXTBOOK_S3_BUCKET')
    try:
        # List objects in the user's metadata folder
        response = s3_client.list_objects_v2(
            Bucket=S3_BUCKET_NAME,
            Prefix=f"metadata/{current_user}/"
        )
        
        books = []
        for obj in response.get('Contents', []):
            metadata = json.loads(s3_client.get_object(Bucket=S3_BUCKET_NAME, Key=obj['Key'])['Body'].read())
            books.append({"title": metadata['title'], "s3_key": metadata['s3_key']})
        
        return books
    except Exception as e:
        logging.error(f"Failed to retrieve user books from S3: {str(e)}")
        raise HTTPException(status_code=500, detail="Failed to retrieve user books")

@app.get("/download-book/{s3_key}")
async def download_book(s3_key: str, current_user: str = Depends(get_current_user)):
    S3_BUCKET_NAME = os.getenv('TEXTBOOK_S3_BUCKET')
    try:
        # Verify that the book belongs to the current user
        if not s3_key.startswith(f"user-uploads/{current_user}/"):
            raise HTTPException(status_code=403, detail="You don't have permission to access this book")
        
        response = s3_client.get_object(Bucket=S3_BUCKET_NAME, Key=s3_key)
        return StreamingResponse(
            response['Body'].iter_chunks(),
            media_type='application/pdf',
            headers={"Content-Disposition": f"attachment; filename={s3_key.split('/')[-1]}"}
        )
    except ClientError as e:
        if e.response['Error']['Code'] == 'NoSuchKey':
            raise HTTPException(status_code=404, detail="Book not found")
        logging.error(f"Failed to download PDF from S3: {str(e)}")
        raise HTTPException(status_code=500, detail="Failed to download PDF")

def generate_game_idea(text: str, learning_profile: str, max_attempts=1, max_tokens=4096):
    prompt = f"""Based on the following materials and the user's learning profile, suggest a simple interactive game idea that reinforces the key concepts. The game should:
    1. Be implementable in JavaScript
    2. Reinforce one or more key concepts from the materials
    3. Be engaging and educational for students
    4. Not be resource intensive and be able to run on a web browser that is using React for frontend
    5. Be tailored to the user's learning style as described in their profile

    Primary Content from Collection:
    {text}

    User's learning profile: {learning_profile}

    Now, provide a game idea that integrates concepts from the available materials."""

    full_response = ""
    for i in range(max_attempts):
        try:
            native_request = {
                'anthropic_version': 'bedrock-2023-05-31',
                'max_tokens': max_tokens,
                'temperature': 0.7,
                'top_p': 0.9,
                'messages': [
                    {
                        'role': 'user',
                        'content': [{'type': 'text', 'text': prompt}],
                    }
                ],
            }

            request = json.dumps(native_request)

            response = bedrock.invoke_model_with_response_stream(
                modelId="us.anthropic.claude-3-5-haiku-20241022-v1:0",
                body=request
            )

            for event in response['body']:
                chunk = json.loads(event['chunk']['bytes'])
                if chunk['type'] == 'content_block_delta':
                    full_response += chunk['delta'].get('text', '')

            if full_response.endswith(".") and len(full_response) < max_tokens * 0.9:
                break

        except ClientError as e:
            logging.error(f"Error calling Bedrock: {e}")
            return f"Error: {str(e)}"

    return full_response

def translate_text(text, target_language):
    try:
        response = translate.translate_text(
            Text=text,
            SourceLanguageCode='en',
            TargetLanguageCode=target_language
        )
        return response['TranslatedText']
    except ClientError as e:
        logging.error(f"Error translating text: {e}")
        return None

@app.post("/generate-narrative/{collection_id}")
async def generate_narrative_endpoint(
    collection_id: str,
    current_user: str = Depends(get_current_user)
):
    try:
        # Get collection content and learning profile
        content = await get_collection_content(collection_id, current_user)
        learning_profile = get_learning_profile(current_user)
        
        # Query knowledge base using existing function
        relevant_info = await query_knowledge_base(format_content_for_prompt(content))
        
        # Build prompt with dynamically formatted content
        narrative_prompt = f"""
        You are LearnAI, a GenAI powered learning assistant that adjusts content to the user's learning profile. 
        Generate an extensive, in-depth summary for the following materials, making sure to cover all key concepts 
        while incorporating the provided relevant information and tailoring it to the user's learning profile:

        Primary Content from Collection:
        {format_content_for_prompt(content)}

        Additional Relevant Information from Knowledge Base:
        {relevant_info}

        Learning profile: {learning_profile}

        Please create a comprehensive, detailed walkthrough that:
        1. Thoroughly explains all key concepts from all materials
        2. Integrates information across different content types
        3. Provides multiple examples and applications
        4. Draws connections between different sources
        5. Uses rich analogies and real-world examples
        6. Addresses common misconceptions
        7. Includes thought-provoking questions
        8. Adjusts content to cater to the user's learning profile
        """

        # Generate narrative
        narrative = generate_narrative(narrative_prompt)

        # Generate game idea
        game_response = generate_game_idea(format_content_for_prompt(content), learning_profile)
        
        # Generate game code
        game_code_response = await generate_game_code(GameIdeaRequest(game_idea=game_response))
        game_code = game_code_response.get("code", "")
        
        # Generate diagrams
        diagrams = await generate_diagrams(format_content_for_prompt(content), narrative, learning_profile)

        return {
            "narrative": narrative,
            "game_idea": game_response,
            "game_code": game_code,
            "diagrams": diagrams
        }

    except Exception as e:
        logging.error(f"Error generating narrative: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to generate narrative: {str(e)}"
        )

@app.post("/transcribe")
async def transcribe_audio(audio: UploadFile = File(...)):
    try:
        # Save the uploaded file temporarily
        temp_file_path = f"/tmp/{uuid.uuid4()}.wav"
        with open(temp_file_path, "wb") as buffer:
            buffer.write(await audio.read())
        
        job_name = f"transcribe_job_{int(time.time())}"
        job_uri = f"s3://YOUR_S3_BUCKET/{audio.filename}"
        
        # Upload the file to S3 (you need to implement this)
        # upload_to_s3(temp_file_path, "YOUR_S3_BUCKET", audio.filename)
        
        transcribe.start_transcription_job(
            TranscriptionJobName=job_name,
            Media={'MediaFileUri': job_uri},
            MediaFormat='wav',
            LanguageCode='en-US'
        )
        
        while True:
            status = transcribe.get_transcription_job(TranscriptionJobName=job_name)
            if status['TranscriptionJob']['TranscriptionJobStatus'] in ['COMPLETED', 'FAILED']:
                break
            time.sleep(5)
        
        if status['TranscriptionJob']['TranscriptionJobStatus'] == 'COMPLETED':
            result = transcribe.get_transcription_job(TranscriptionJobName=job_name)
            transcript_uri = result['TranscriptionJob']['Transcript']['TranscriptFileUri']
            # Fetch and parse the transcript JSON from the URI
            # You'll need to implement this part
            transcript_text = fetch_and_parse_transcript(transcript_uri)
            return {"transcript": transcript_text}
        else:
            return {"error": "Transcription failed"}

    except Exception as e:
        logging.exception("Error in transcribe_audio")
        raise HTTPException(status_code=500, detail=str(e))

def get_learning_profile(user_id):
    if not user_id:
        logging.error("get_learning_profile called with empty user_id")
        return "Learning profile not available."

    try:
        profile_key = f"learning_profiles/{user_id}.json"
        response = s3_client.get_object(Bucket=S3_BUCKET_NAME, Key=profile_key)
        profile_data = json.loads(response['Body'].read().decode('utf-8'))
        return profile_data.get('description', 'Learning profile not available.')
    except ClientError as e:
        if e.response['Error']['Code'] == 'NoSuchKey':
            logging.error(f"Learning profile not found for user_id: {user_id}")
        else:
            logging.error(f"Error fetching learning profile: {str(e)}")
        return "Learning profile not available."
    except Exception as e:
        logging.error(f"Unexpected error in get_learning_profile: {str(e)}")
        return "Learning profile not available."
    
def generate_narrative(prompt: str, max_attempts=1, max_tokens=8192):
    full_response = ""
    for i in range(max_attempts):
        try:
            native_request = {
                'anthropic_version': 'bedrock-2023-05-31',
                'max_tokens': max_tokens,
                'temperature': 0.7,
                'top_p': 0.9,
                'messages': [
                    {
                        'role': 'user',
                        'content': [{'type': 'text', 'text': prompt}],
                    }
                ],
            }

            request = json.dumps(native_request)

            response = bedrock.invoke_model_with_response_stream(
                modelId="us.anthropic.claude-3-5-haiku-20241022-v1:0",
                body=request
            )

            for event in response['body']:
                chunk = json.loads(event['chunk']['bytes'])
                if chunk['type'] == 'content_block_delta':
                    full_response += chunk['delta'].get('text', '')

            if full_response.endswith(".") and len(full_response) < max_tokens * 0.9:
                break

        except ClientError as e:
            logging.error(f"Error calling Bedrock: {e}")
            return f"Error: {str(e)}"

    return full_response.strip()

class GameIdeaRequest(BaseModel):
    game_idea: str

@app.post("/generate-game-code")
async def generate_game_code(request: GameIdeaRequest):
    game_idea = request.game_idea
    
    prompt = f"""Create a fully functional React component for the following game idea that integrates concepts from multiple learning materials:

    {game_idea}

    The component will be rendered within a DynamicGameComponent. Your task is to generate the code that will be passed to this component.

            ```javascript
            const DynamicGameComponent = (((this is where the game code will be passed in as an argument))) => {{
              const [error, setError] = useState(null);
              const [GameComponent, setGameComponent] = useState(null);

              useEffect(() => {{
                setError(null);
                if (gameCode.startsWith("Error generating game code:")) {{
                  setError(gameCode);
                  return;
                }}
                try {{
                  // Your generated code will be inserted here
                  const ComponentFunction = new Function('React', 'useState', 'useEffect', 'MathJax', `
                    return function Game() {{
                      // Your generated code goes here
                    }}
                  `);

                  const CreatedComponent = () => {{
                    return (
                      <ErrorBoundary>
                        {{ComponentFunction(React, React.useState, React.useEffect, MathJax)}}
                      </ErrorBoundary>
                    );
                  }};
                  setGameComponent(() => CreatedComponent);
                }} catch (err) {{
                  console.error('Error creating game component:', err);
                  setError('Error: (((this is where error output would go))));
                }}
              }}, [gameCode]);

              // ... rest of the component
            }};
            ```

            Requirements:
                1. Use React hooks (useState, useEffect, useRef, useCallback) without React. prefix
                2. Use React.createElement for all element creation (no JSX)
                3. Return a single root element (usually a div) containing all other elements 
                4. Ensure all variables and functions are properly declared
                5. Do not use any external libraries or components not provided
                6. Provide ONLY the JavaScript code, without any explanations or markdown formatting
                7. Do not include 'return function Game() {{' at the beginning or '}}' at the end
                8. Use proper JavaScript syntax (no semicolons after blocks or object literals in arrays)
                9. Do not use 'function' as a variable name, as it is a reserved keyword in JavaScript. Use 'func' or 'mathFunction' instead
                10. Create instructions for the user on how to play the game in the game component and how it relates to the chapter content
                11. When evaluating mathematical expressions or functions, use a safe evaluation method instead of 'eval'. For example:
                    - For simple arithmetic, use basic JavaScript operations
                    - For more complex functions, define them explicitly (e.g., Math.sin, Math.cos, etc.)
                12. Ensure all variables used in calculations are properly defined and initialized
                13. Use try-catch blocks when performing calculations to handle potential errors gracefully
                14. For keyboard input:
                    - Use the useEffect hook to add and remove event listeners for keyboard events
                    - In the event listener, call e.preventDefault() to prevent default browser behavior (like scrolling)
                    - Focus on a game element (like the canvas) when the component mounts to ensure it captures keyboard events
                15. Add a button to start/restart the game, and only capture keyboard input when the game is active
                16. Ensure that the current equation is always visible and properly rendered using plain text or another method
                17. To prevent scrolling when using arrow keys:
                    - Add 'tabIndex={0}' to the game container div to make it focusable
                    - In the useEffect for keyboard events, check if the game container has focus before handling key presses
                18. Display the current function prominently using plain text, and update it whenever it changes
                19. Use requestAnimationFrame for the game loop to ensure smooth animation
                20. Add error checking before accessing array elements or object properties
                21. Use optional chaining (?.) when accessing nested properties to prevent errors
                22. The background color of the dynamic game component is white keep this as the background color of the game.
                23. Do not forget to include instructions for the user on how to play the game in the game component and how it relates to the chapter content as text in the game component.
                24. Container Sizing Requirements:
                    - The container is the parent of the game component and is in a portrait orientation occupying vertically the right half of the screen(it's aspect ratio is 5:6)
                    - The game must automatically scale to fit its container width without scrollbars
                    - Use relative units (%, vh, vw) instead of fixed pixel values
                    - All game elements must be sized relative to their container
                    - The game container uses width: 100% and height: 100%
                    - Add a useEffect hook to handle window resizing and maintain proper scaling
                    - Ensure the game's aspect ratio is maintained while fitting the width
                    - Use transform: scale() if needed to ensure proper fitting
                    - Set initial dimensions using percentages of the parent container
                    - Listen for container size changes and update game element sizes accordingly
                    - Use getBoundingClientRect() to get accurate container dimensions
                    - Apply CSS transform-origin: top left when scaling
                                                                
            Generate the game code now, remember to not include any explanations or comments, just the code:"""

    try:
        body = json.dumps({
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 3000,
            "messages": [
                {
                    "role": "user",
                    "content": [{"type": "text", "text": prompt}]
                }
            ],
            "temperature": 0.7,
            "top_p": 0.9,
        })

        response = bedrock.invoke_model_with_response_stream(
            modelId="us.anthropic.claude-3-5-sonnet-20241022-v2:0",
            body=body
        )

        generated_code = ""
        for event in response['body']:
            chunk = json.loads(event['chunk']['bytes'])
            if chunk['type'] == 'content_block_delta':
                generated_code += chunk['delta'].get('text', '')

        logging.debug(f"Generated game code:\n{generated_code}")

        return {"code": generated_code.strip()}

    except Exception as e:
        logging.error(f"Error generating game code: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error generating game code: {str(e)}")

def validate_js_syntax(code):
    try:
        esprima.parseScript(code)
        return True
    except esprima.Error as e:
        logging.error(f"Syntax error in generated code: {e}")
        return False

def post_process_game_code(code):
    # Remove any potential wrapper function or extra brackets
    code = re.sub(r'^return\s*function\s*\w*\s*\(\)\s*{', '', code)
    code = re.sub(r'}\s*$', '', code)
    
    # Remove semicolons after object literals in array definitions
    code = re.sub(r'({[^}]+});(?=\s*[}\]])', r'\1', code)
    
    # Fix arrow function syntax
    code = re.sub(r'(\w+)\s*=>\s*{', r'\1 => {', code)
    
    # Remove semicolons after function/if/else blocks
    code = re.sub(r'(}\s*);(\s*else|\s*[)\]])', r'\1\2', code)
    
    # Ensure all state variables are properly declared
    state_vars = re.findall(r'const \[(\w+), set\w+\] = useState\((.*?)\);', code)
    for var, initial_value in state_vars:
        if initial_value.strip() == '':
            code = code.replace(f'const [{var}, set{var.capitalize()}] = useState();', 
                                f'const [{var}, set{var.capitalize()}] = useState(null);')
    
    # Fix MathJax.Node syntax if necessary
    code = code.replace('React.createElement(MathJax.Node,', 'React.createElement(MathJax,')
    
    # Remove semicolons inside object literals
    code = re.sub(r'({[^}]+};\s*}', r'\1 }', code)
    
    return code.strip()

@app.post("/api/chat")
async def chat(request: Request):
    try:
        data = await request.json()
        user_message = data.get('message')
        user_id = data.get('userId')
        file_id = data.get('fileId')
        section_name = data.get('sectionName')
        language = data.get('language', 'en')
        force_regenerate = data.get('forceRegenerate')
        logging.info(f"Received chat request for section: {section_name}, force_regenerate: {force_regenerate}")

        if not user_message or not user_id or not file_id or not section_name:
            raise HTTPException(status_code=400, detail="Missing required parameters")

        # Retrieve the extracted text from S3
        extracted_text_key = f"extracted-text/{user_id}/{file_id}/section_{section_name}.txt"
        try:
            extracted_text_obj = s3_client.get_object(Bucket=S3_BUCKET_NAME, Key=extracted_text_key)
            extracted_text = extracted_text_obj['Body'].read().decode('utf-8')
        except ClientError as e:
            logging.error(f"Error retrieving extracted text: {e}")
            extracted_text = "No extracted text available."

        # Retrieve the generated summary narrative from S3
        narrative_key = f"narratives/{user_id}/{file_id}/{section_name}_{force_regenerate}.json"
        try:
            narrative_obj = s3_client.get_object(Bucket=S3_BUCKET_NAME, Key=narrative_key)
            narrative_data = json.loads(narrative_obj['Body'].read().decode('utf-8'))
            generated_summary = narrative_data.get('narrative', '')
        except ClientError as e:
            logging.error(f"Error retrieving narrative: {e}")
            generated_summary = "No generated summary available."

        # Retrieve chat history from S3
        chat_history_key = f"chat-history/{user_id}/{file_id}/{section_name}.json"
        try:
            chat_history_obj = s3_client.get_object(Bucket=S3_BUCKET_NAME, Key=chat_history_key)
            chat_history = json.loads(chat_history_obj['Body'].read().decode('utf-8'))
        except ClientError as e:
            if e.response['Error']['Code'] == 'NoSuchKey':
                chat_history = []
            else:
                logging.error(f"Error retrieving chat history: {e}")
                chat_history = []

        # Retrieve knowledge base information
        knowledge_base_info = await query_knowledge_base(extracted_text[:2000])

        if knowledge_base_info is None:
            knowledge_base_info = ""
            logging.error("Knowledge base info is None")
            
        # Translate user message to English if not already in English
        if language != 'en':
            user_message = translate_text(user_message, 'en')

        context = f"""You are LearnAI, a GenAI powered learning assistant that adjusts textbook content to the user's learning profile. 
        The current section content is: {extracted_text}... 
        The generated summary of this section is: {generated_summary}...
        Relevant information from the knowledge base is: {knowledge_base_info[:2000]}...
        Please ensure your responses are relevant to this topic."""

        prompt = f"""Here's the user's learning profile: {get_learning_profile(user_id)}

        {context}

        Remember the context of the previous messages in this conversation. Here's the student's latest question:

        {user_message}

        Provide a helpful, accurate, and concise answer based on the given context, your general knowledge, and the conversation history. Make sure to reference the section content in your answer. Answer the question but be as concise as possible (4-6 sentences)."""

        # Generate AI response using the chat function
        ai_response = generate_chat_response(prompt, chat_history)

        # Translate AI response if not in English
        if language != 'en':
            ai_response = translate_text(ai_response, language)

        # Update chat history
        chat_history.append({"user": "You", "text": data.get('message')})  # Use original message for history
        chat_history.append({"user": "AI", "text": ai_response})

        # Save updated chat history to S3
        s3_client.put_object(
            Bucket=S3_BUCKET_NAME,
            Key=chat_history_key,
            Body=json.dumps(chat_history),
            ContentType='application/json'
        )

        return {"reply": ai_response}

    except Exception as e:
        logging.exception("Error in chat endpoint")
        raise HTTPException(status_code=500, detail=str(e))
    
def generate_chat_response(prompt: str, chat_history: list, max_tokens: int = 500, max_history_tokens: int = 1000) -> str:
    full_response = ""
    
    try:
        # Prepare the chat history for the model
        formatted_history = []
        total_tokens = 0
        for msg in reversed(chat_history):
            role = "user" if msg['user'] == 'You' else "assistant"
            content = msg['text']
            message_tokens = len(content.split())  # Simple token count estimation
            if total_tokens + message_tokens > max_history_tokens:
                break
            formatted_history.insert(0, {"role": role, "content": content})
            total_tokens += message_tokens

        # Add the current prompt
        formatted_history.append({"role": "user", "content": prompt})

        native_request = {
            'anthropic_version': 'bedrock-2023-05-31',
            'max_tokens': max_tokens,
            'messages': formatted_history,
            'temperature': 0.7,
            'top_k': 250,
            'top_p': 1,
            'stop_sequences': ['\n\nHuman:'],
        }

        request = json.dumps(native_request)

        response = bedrock.invoke_model_with_response_stream(
            modelId="us.anthropic.claude-3-5-haiku-20241022-v1:0",
            body=request
        )

        for event in response['body']:
            chunk = json.loads(event['chunk']['bytes'])
            if chunk['type'] == 'content_block_delta':
                full_response += chunk['delta'].get('text', '')

        return full_response.strip()

    except ClientError as e:
        logging.error(f"Error calling Bedrock: {e}")
        return f"Error: {str(e)}"

@app.post("/transcribe")
async def transcribe_audio(audio: UploadFile = File(...)):
    # Save the uploaded file temporarily
    temp_file_path = f"/tmp/{uuid.uuid4()}.wav"
    with open(temp_file_path, "wb") as buffer:
        buffer.write(await audio.read())
    
    job_name = f"transcribe_job_{int(time.time())}"
    job_uri = f"s3://YOUR_S3_BUCKET/{audio.filename}"
    
    # Upload the file to S3 (you need to implement this)
    # upload_to_s3(temp_file_path, "YOUR_S3_BUCKET", audio.filename)
    
    transcribe.start_transcription_job(
        TranscriptionJobName=job_name,
        Media={'MediaFileUri': job_uri},
        MediaFormat='wav',
        LanguageCode='en-US'
    )
    
    while True:
        status = transcribe.get_transcription_job(TranscriptionJobName=job_name)
        if status['TranscriptionJob']['TranscriptionJobStatus'] in ['COMPLETED', 'FAILED']:
            break
        time.sleep(5)
    
    if status['TranscriptionJob']['TranscriptionJobStatus'] == 'COMPLETED':
        result = transcribe.get_transcription_job(TranscriptionJobName=job_name)
        transcript_uri = result['TranscriptionJob']['Transcript']['TranscriptFileUri']
        # Fetch and parse the transcript JSON from the URI
        # You'll need to implement this part
        transcript_text = fetch_and_parse_transcript(transcript_uri)
        return {"transcript": transcript_text}
    else:
        return {"error": "Transcription failed"}

@app.post("/api/synthesize-speech")
async def synthesize_speech(request: Request):
    data = await request.json()
    text = data.get('text')
    language = data.get('language', 'en-US')

    try:
        # Map language codes to Polly voice IDs
        voice_map = {
            'en-US': 'Joanna',
            'es-ES': 'Conchita',
            'fr-FR': 'Celine',
            'de-DE': 'Marlene',
            # Add more languages and voices as needed
        }

        response = polly_client.synthesize_speech(
            Text=text,
            OutputFormat='mp3',
            VoiceId=voice_map.get(language, 'Joanna')
        )

        if "AudioStream" in response:
            with closing(response["AudioStream"]) as stream:
                output = os.path.join(gettempdir(), "speech.mp3")
                try:
                    with open(output, "wb") as file:
                        file.write(stream.read())
                except IOError as error:
                    print(error)
                    raise HTTPException(status_code=500, detail="Error writing audio stream to file")

            return FileResponse(output, media_type='audio/mpeg', filename="speech.mp3")
        else:
            raise HTTPException(status_code=500, detail="Could not stream audio")

    except (BotoCoreError, ClientError) as error:
        print(error)
        raise HTTPException(status_code=500, detail=str(error))

@app.post("/translate")
async def translate_text_endpoint(request: Request):
    data = await request.json()
    text = data.get('text')
    target_language = data.get('target_language')

    if not text or not target_language:
        raise HTTPException(status_code=400, detail="Text and target language are required")

    translated_text = translate_text(text, target_language)

    if translated_text is None:
        raise HTTPException(status_code=500, detail="Translation failed")

    return {"translated_text": translated_text}

@app.get("/user-textbooks")
async def get_user_textbooks(current_user: str = Depends(get_current_user)):
    try:
        # List objects in the user's metadata folder
        response = s3_client.list_objects_v2(
            Bucket=S3_BUCKET_NAME,
            Prefix=f"metadata/{current_user}/"
        )
        
        textbooks = []
        for obj in response.get('Contents', []):
            metadata = json.loads(s3_client.get_object(Bucket=S3_BUCKET_NAME, Key=obj['Key'])['Body'].read())
            if metadata.get('document_type') == 'textbook':
                textbooks.append({
                    "title": metadata['title'],
                    "s3_key": metadata['s3_key'],
                    "upload_date": obj['LastModified'].strftime("%Y-%m-%d %H:%M:%S")
                })
        
        return textbooks
    except Exception as e:
        logging.error(f"Failed to retrieve user textbooks from S3: {str(e)}")
        raise HTTPException(status_code=500, detail="Failed to retrieve user textbooks")

@app.get("/textbook-structure/{user_id}/{file_id}/{filename}")
async def get_textbook_structure(
    user_id: str = Path(...),
    file_id: str = Path(...),
    filename: str = Path(...),
    current_user: str = Depends(get_current_user)
):
    if user_id != current_user:
        raise HTTPException(status_code=403, detail="You don't have permission to access this file")

    metadata_key = f"metadata/{user_id}/{file_id}_{filename}.json"
    print(f"Constructed metadata key: {metadata_key}")

    try:
        # Fetch the metadata file from S3
        response = s3_client.get_object(Bucket=S3_BUCKET_NAME, Key=metadata_key)
        metadata = json.loads(response['Body'].read().decode('utf-8'))
        print(f"Metadata content: {json.dumps(metadata, indent=2)}")
        print(f"Successfully fetched metadata for key: {metadata_key}")

        # Download and save the PDF
        temp_path = download_and_save_pdf(user_id, file_id, filename)

        # Extract the table of contents from the metadata
        toc = metadata.get('table_of_contents', [])
        print(f"Extracted table of contents: {toc}")

        # Transform the table of contents into the required structure
        book_structure = {
            "chapters": []
        }

        for chapter in toc:
            chapter_structure = {
                "id": chapter['number'],
                "title": f"{chapter['number']}: {chapter['title']}",
                "sections": []
            }
            for section in chapter.get('sections', []):
                chapter_structure['sections'].append({
                    "id": section['title'].split()[0],
                    "title": section['title']
                })
            book_structure['chapters'].append(chapter_structure)

        print(f"Transformed book structure: {book_structure}")
        return book_structure

    except ClientError as e:
        if e.response['Error']['Code'] == 'NoSuchKey':
            print(f"Metadata file not found. Key: {metadata_key}")
            raise HTTPException(status_code=404, detail="Textbook structure not found")
        print(f"Error fetching metadata from S3: {str(e)}")
        raise
    except Exception as e:
        print(f"Unexpected error in get_textbook_structure: {str(e)}")
        raise HTTPException(status_code=500, detail="An unexpected error occurred")

@app.get("/list-metadata/{user_id}")
async def list_metadata(user_id: str, current_user: str = Depends(get_current_user)):
    if user_id != current_user:
        raise HTTPException(status_code=403, detail="You don't have permission to access this data")
    
    try:
        response = s3_client.list_objects_v2(
            Bucket=S3_BUCKET_NAME,
            Prefix=f"metadata/{user_id}/"
        )
        
        files = [obj['Key'] for obj in response.get('Contents', [])]
        return {"metadata_files": files}
    except Exception as e:
        logging.error(f"Error listing metadata files: {str(e)}")
        raise HTTPException(status_code=500, detail="Failed to list metadata files")

# Create a cache with a maximum of 100 items and a 1-hour TTL
section_cache = TTLCache(maxsize=100, ttl=3600)

@app.get("/get-section-pdf/{user_id}/{file_id}/{filename}/{section_id}")
async def get_section_pdf(user_id: str, file_id: str, filename: str, section_id: str, current_user: str = Depends(get_current_user)):
    cache_key = f"{user_id}_{file_id}_{filename}_{section_id}"
    
    # Check if the section is in the cache
    if cache_key in section_cache:
        logging.info(f"Returning cached section for {cache_key}")
        return StreamingResponse(io.BytesIO(section_cache[cache_key]), media_type="application/pdf", 
                                 headers={"Content-Disposition": f"attachment; filename={section_id}.pdf"})

    logging.info(f"Fetching section PDF for file_id: {file_id}, filename: {filename}, section_id: {section_id}")
    
    # Fetch metadata (same as before)
    metadata_key = f"metadata/{user_id}/{file_id}_{filename}.json"
    try:
        metadata_obj = s3_client.get_object(Bucket=S3_BUCKET_NAME, Key=metadata_key)
        metadata = json.loads(metadata_obj['Body'].read().decode('utf-8'))
    except Exception as e:
        logging.error(f"Error fetching metadata: {str(e)}")
        raise HTTPException(status_code=404, detail="Metadata not found")

    # Find the section in the table of contents (same as before)
    toc = metadata.get('table_of_contents', [])
    section = None
    next_section = None
    for chapter in toc:
        for i, s in enumerate(chapter['sections']):
            if s['title'].startswith(section_id):
                section = s
                if i + 1 < len(chapter['sections']):
                    next_section = chapter['sections'][i + 1]
                break
        if section:
            break

    if not section:
        raise HTTPException(status_code=404, detail="Section not found")

    # Create a sanitized filename from the section title
    sanitized_title = "".join(c for c in section['title'] if c.isalnum() or c in (' ', '-', '_')).rstrip()
    pdf_filename = f"{sanitized_title}.pdf"
    print(f"Sanitized title: {sanitized_title}, PDF filename: {pdf_filename}")

    # Determine the end page (same as before)
    if next_section:
        end_page = next_section['page'] - 1
    else:
        chapter_index = toc.index(chapter)
        if chapter_index + 1 < len(toc):
            end_page = toc[chapter_index + 1]['sections'][0]['page'] - 1
        else:
            end_page = section['page'] + 50  # Arbitrary number, adjust as needed

    logging.info(f"Extracting pages {section['page']} to {end_page}")

    # Use the existing temporary file
    temp_key = f"{user_id}_{file_id}_{filename}"
    if temp_key not in temp_pdfs:
        # If the temporary file doesn't exist, download it
        temp_path = download_and_save_pdf(user_id, file_id, filename)
    else:
        temp_path = temp_pdfs[temp_key]

    # Extract the section pages using PyMuPDF
    pdf_document = fitz.open(temp_path)
    output_pdf = fitz.open()

    for page_num in range(section['page'] - 1, min(end_page, len(pdf_document))):
        output_pdf.insert_pdf(pdf_document, from_page=page_num, to_page=page_num)

    # Save the extracted pages to a new PDF
    output_buffer = io.BytesIO()
    output_pdf.save(output_buffer)
    output_buffer.seek(0)

    # Cache the extracted section
    section_cache[cache_key] = output_buffer.getvalue()

    logging.info("Successfully extracted and prepared PDF section")

    return StreamingResponse(output_buffer, media_type="application/pdf", 
                             headers={"Content-Disposition": f"attachment; filename={pdf_filename}"})

# Global dictionary to store temporary file paths
temp_pdfs = {}

def download_and_save_pdf(user_id, file_id, filename):
    temp_key = f"{user_id}_{file_id}_{filename}"
    if temp_key in temp_pdfs:
        logging.info(f"Using existing temporary PDF: {temp_pdfs[temp_key]}")
        return temp_pdfs[temp_key]

    s3_key = f"user-uploads/{user_id}/{file_id}_{filename}"
    temp_dir = tempfile.gettempdir()
    temp_path = os.path.join(temp_dir, f"{file_id}_{filename}")
    
    try:
        response = s3_client.get_object(Bucket=S3_BUCKET_NAME, Key=s3_key)
        with open(temp_path, 'wb') as f:
            f.write(response['Body'].read())
        temp_pdfs[temp_key] = temp_path
        logging.info(f"PDF downloaded and saved temporarily: {temp_path}")
        return temp_path
    except Exception as e:
        logging.error(f"Error downloading PDF: {str(e)}")
        raise HTTPException(status_code=500, detail="Failed to download PDF")

# Function to clean up temporary files
def cleanup_temp_files():
    for temp_path in temp_pdfs.values():
        try:
            os.remove(temp_path)
            logging.info(f"Temporary file removed: {temp_path}")
        except Exception as e:
            logging.error(f"Error removing temporary file {temp_path}: {str(e)}")

# Register the cleanup function to run when the server shuts down
atexit.register(cleanup_temp_files)

@app.post("/process-pdf-section")
async def process_pdf_section(
    user_id: str = Form(...),
    file_id: str = Form(...),
    filename: str = Form(...),
    section_name: str = Form(...),
    current_user: str = Depends(get_current_user)
):
    try:
        logging.info(f"Received request for process_pdf_section:")
        logging.info(f"user_id: {user_id}")
        logging.info(f"file_id: {file_id}")
        logging.info(f"filename: {filename}")
        logging.info(f"section_name: {section_name}")
        logging.info(f"current_user: {current_user}")

        if user_id != current_user:
            logging.warning(f"Permission denied: user_id ({user_id}) does not match current_user ({current_user})")
            raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="You don't have permission to access this file")

        # Fetch metadata
        metadata_key = f"metadata/{user_id}/{file_id}_{filename}.json"
        try:
            metadata_obj = s3_client.get_object(Bucket=S3_BUCKET_NAME, Key=metadata_key)
            metadata = json.loads(metadata_obj['Body'].read().decode('utf-8'))
        except Exception as e:
            logging.error(f"Error fetching metadata: {str(e)}")
            raise HTTPException(status_code=404, detail="Metadata not found")

        # Find the section in the table of contents
        toc = metadata.get('table_of_contents', [])
        section = None
        next_section = None
        for chapter in toc:
            for i, s in enumerate(chapter['sections']):
                if s['title'] == section_name:
                    section = s
                    if i + 1 < len(chapter['sections']):
                        next_section = chapter['sections'][i + 1]
                    break
            if section:
                break

        if not section:
            raise HTTPException(status_code=404, detail="Section not found")

        start_page = section['page']
        if next_section:
            end_page = next_section['page'] - 1
        else:
            chapter_index = toc.index(chapter)
            if chapter_index + 1 < len(toc):
                end_page = toc[chapter_index + 1]['sections'][0]['page'] - 1
            else:
                end_page = start_page + 50  # Arbitrary number, adjust as needed

        logging.info(f"Processing section '{section_name}' from page {start_page} to {end_page}")

        try:
            s3_key = f"user-uploads/{user_id}/{file_id}_{filename}"
            extracted_text_key = f"extracted-text/{user_id}/{file_id}/section_{section_name}.txt"

            # Check if extracted text already exists in S3
            try:
                existing_text = s3_client.get_object(Bucket=S3_BUCKET_NAME, Key=extracted_text_key)
                return {"extracted_text": existing_text['Body'].read().decode('utf-8')}
            except s3_client.exceptions.NoSuchKey:
                # If the file doesn't exist, proceed with extraction
                pass

            extracted_text = ""

            # Download the PDF from S3
            response = s3_client.get_object(Bucket=S3_BUCKET_NAME, Key=s3_key)
            pdf_content = response['Body'].read()

            # Create a temporary file to store the downloaded PDF
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_file:
                temp_file.write(pdf_content)
                temp_path = temp_file.name

            # Process each page individually
            for page_num in range(start_page, end_page + 1):
                logging.info(f"Processing page {page_num}")
                
                # Extract and compress the single page
                with pikepdf.Pdf.open(temp_path) as pdf:
                    new_pdf = pikepdf.Pdf.new()
                    new_pdf.pages.append(pdf.pages[page_num - 1])
                    
                    # Save the compressed page to a BytesIO object
                    output = io.BytesIO()
                    new_pdf.save(output, compress_streams=True, object_stream_mode=pikepdf.ObjectStreamMode.generate)
                    output.seek(0)
                    file_bytes = output.getvalue()

                logging.info(f"Compressed PDF size for page {page_num}: {len(file_bytes)} bytes")
                
                # Process with Textract
                response = textract_client.analyze_document(
                    Document={'Bytes': file_bytes},
                    FeatureTypes=['TABLES', 'FORMS']
                )
                
                # Extract text and table data
                page_text = extract_text_and_tables(response)
                extracted_text += f"Page {page_num}:\n{page_text}\n\n"

            # Save the extracted text to S3
            s3_client.put_object(Bucket=S3_BUCKET_NAME, Key=extracted_text_key, Body=extracted_text)

            return {"extracted_text": extracted_text}

        except Exception as e:
            logging.error(f"Error processing PDF section: {str(e)}")
            raise HTTPException(status_code=500, detail="Error processing PDF section")
        finally:
            # Clean up the temporary file
            if 'temp_path' in locals():
                os.unlink(temp_path)

    except ValueError as ve:
        logging.error(f"ValueError in process_pdf_section: {str(ve)}")
        raise HTTPException(status_code=status.HTTP_422_UNPROCESSABLE_ENTITY, detail=str(ve))
    except Exception as e:
        logging.error(f"Unexpected error in process_pdf_section: {str(e)}")
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail="An unexpected error occurred")

def extract_text_and_tables(textract_response):
    extracted_content = ""
    for item in textract_response['Blocks']:
        if item['BlockType'] == 'LINE':
            extracted_content += item['Text'] + "\n"
        elif item['BlockType'] == 'TABLE':
            extracted_content += "Table:\n"
            for cell in item['Relationships'][0]['Ids']:
                cell_block = next(block for block in textract_response['Blocks'] if block['Id'] == cell)
                if 'Text' in cell_block:
                    extracted_content += cell_block['Text'] + "\t"
                extracted_content += "\n"
            extracted_content += "\n"
    return extracted_content

async def generate_diagrams(content: str, narrative: str, learning_profile: str):
    prompt = f"""Based on the following materials, generated summary, and the user's learning profile, create a set of diagrams that illustrate the key concepts:

    Primary Content from Collection:
    {content}

    Generated Summary:
    {narrative}

    User's Learning Profile:
    {learning_profile}

    Please create diagrams that:
    1. Illustrate the main concepts and their relationships
    2. Are clear and easy to understand
    3. Are tailored to the user's learning style as described in their profile
    4. Use appropriate visual representations (e.g., flowcharts, mind maps, etc.)

    Strict guidelines for Mermaid syntax:
    - Start each diagram with 'graph TD' on its own line
    - Use only alphanumeric characters and underscores for node IDs (e.g., A, B, C, or Node1, Node2, Node3)
    - Use square brackets for node labels: [Label text]
    - Use only --> for arrows (no labels or other arrow types)
    - Put each node and connection on its own line
    - Do not use any special characters, mathematical symbols, or subscripts in labels
    - If referring to mathematical concepts, use words instead of symbols (e.g., "First Derivative" instead of "f'(x)")
    - Ensure all nodes are connected in a logical flow
    - Keep labels short and concise
    - Limit each diagram to a maximum of 10 nodes for clarity

    Example of correct Mermaid syntax:
    ```mermaid
    graph TD
        A[First Concept] --> B[Second Concept]
        B --> C[Third Concept]
        C --> D[Fourth Concept]
        B --> E[Fifth Concept]
    ```

    Provide 2-3 diagrams in correct Mermaid syntax, each enclosed in ```mermaid and ``` tags."""

    try:
        native_request = {
            'anthropic_version': 'bedrock-2023-05-31',
            'max_tokens': 4096,
            'temperature': 0.7,
            'top_p': 0.9,
            'messages': [
                {
                    'role': 'user',
                    'content': [{'type': 'text', 'text': prompt}],
                }
            ],
        }

        request = json.dumps(native_request)

        response = bedrock.invoke_model_with_response_stream(
            modelId="us.anthropic.claude-3-5-haiku-20241022-v1:0",
            body=request
        )

        full_response = ""
        for event in response['body']:
            chunk = json.loads(event['chunk']['bytes'])
            if chunk['type'] == 'content_block_delta':
                full_response += chunk['delta'].get('text', '')

        # Extract Mermaid diagrams from the response
        mermaid_diagrams = re.findall(r'```mermaid\n(.*?)\n```', full_response, re.DOTALL)
        
        # Post-process each diagram
        processed_diagrams = [post_process_mermaid(diagram) for diagram in mermaid_diagrams]

        return processed_diagrams

    except Exception as e:
        logging.error(f"Error generating diagrams: {str(e)}")
        return []

def post_process_mermaid(diagram):
    # Remove any empty lines
    lines = [line for line in diagram.split('\n') if line.strip()]
    
    # Ensure the diagram starts with 'graph TD'
    if lines[0].strip() != 'graph TD':
        lines.insert(0, 'graph TD')
    
    # Process each line
    processed_lines = []
    for line in lines:
        # Remove any special characters from node IDs
        line = re.sub(r'([A-Za-z0-9_]+)', lambda m: re.sub(r'[^A-Za-z0-9_]', '', m.group(1)), line)
        
        # Ensure proper formatting for node labels
        line = re.sub(r'\[(.*?)\]', lambda m: f"[{m.group(1).replace('_', ' ')}]", line)
        
        # Ensure proper arrow syntax
        line = re.sub(r'-->', ' --> ', line)
        
        processed_lines.append(line)
    
    # Ensure proper indentation
    indented_lines = ['graph TD']
    for line in processed_lines[1:]:
        indented_lines.append('    ' + line.strip())
    
    return '\n'.join(indented_lines)

@app.post("/transcribe-youtube")
async def transcribe_youtube(
    video_url: str = Body(...),
    current_user: str = Depends(get_current_user)
):
    try:
        # Get video info using yt-dlp
        ydl_opts = {
            'quiet': True,
            'no_warnings': True,
            'extract_flat': True
        }
        with YoutubeDL(ydl_opts) as ydl:
            video_info = ydl.extract_info(video_url, download=False)
            video_title = video_info.get('title', 'Untitled Video')
            video_id = video_info.get('id', '')

        # Generate unique IDs for the files
        job_id = str(uuid.uuid4())
        temp_audio_path = f"/tmp/{job_id}.mp3"
        s3_audio_key = f"temp-audio/{job_id}.mp3"

        # Configure yt-dlp options
        ydl_opts = {
            'format': 'bestaudio/best',
            'postprocessors': [{
                'key': 'FFmpegExtractAudio',
                'preferredcodec': 'mp3',
                'preferredquality': '192',
            }],
            'outtmpl': temp_audio_path[:-4],  # Remove .mp3 extension as yt-dlp adds it
            'quiet': True,
            'no_warnings': True
        }

        # Download audio using yt-dlp
        logging.info(f"Attempting to download audio from: {video_url}")
        try:
            with YoutubeDL(ydl_opts) as ydl:
                info = ydl.extract_info(video_url, download=True)
                video_title = info.get('title', 'untitled')
                video_id = info.get('id', 'unknown')
        except Exception as youtube_error:
            logging.error(f"YouTube download error: {str(youtube_error)}")
            raise HTTPException(
                status_code=400,
                detail=f"Failed to download YouTube video: {str(youtube_error)}"
            )

        # Upload audio to S3
        s3_client.upload_file(f"{temp_audio_path[:-4]}.mp3", S3_BUCKET_NAME, s3_audio_key)
        s3_uri = f"s3://{S3_BUCKET_NAME}/{s3_audio_key}"
        
        # Start transcription job
        transcribe.start_transcription_job(
            TranscriptionJobName=job_id,
            Media={'MediaFileUri': s3_uri},
            MediaFormat='mp3',
            LanguageCode='en-US'
        )
        
        # Wait for transcription to complete
        while True:
            status = transcribe.get_transcription_job(TranscriptionJobName=job_id)
            if status['TranscriptionJob']['TranscriptionJobStatus'] in ['COMPLETED', 'FAILED']:
                break
            await asyncio.sleep(5)
        
        # Clean up temporary files
        if os.path.exists(f"{temp_audio_path[:-4]}.mp3"):
            os.remove(f"{temp_audio_path[:-4]}.mp3")
        s3_client.delete_object(Bucket=S3_BUCKET_NAME, Key=s3_audio_key)
        
        if status['TranscriptionJob']['TranscriptionJobStatus'] == 'COMPLETED':
            transcript_uri = status['TranscriptionJob']['Transcript']['TranscriptFileUri']
            
            # Get the transcript JSON
            async with httpx.AsyncClient() as client:
                response = await client.get(transcript_uri)
                transcript_data = response.json()
            
            # Create metadata and transcript objects for S3
            metadata = {
                "video_id": video_id,
                "title": video_title,
                "video_url": video_url,
                "transcription_date": datetime.now().isoformat(),
                "job_id": job_id,
                "original_filename": video_title  # Use video title for collection naming
            }
            
            transcript_text = transcript_data['results']['transcripts'][0]['transcript']
            
            # Save metadata to S3
            metadata_key = f"transcriptions/{current_user}/metadata/{job_id}.json"
            s3_client.put_object(
                Bucket=S3_BUCKET_NAME,
                Key=metadata_key,
                Body=json.dumps(metadata),
                ContentType='application/json'
            )
            
            # Save transcript to S3
            transcript_key = f"transcriptions/{current_user}/content/{job_id}.txt"
            s3_client.put_object(
                Bucket=S3_BUCKET_NAME,
                Key=transcript_key,
                Body=transcript_text,
                ContentType='text/plain'
            )
            
            logging.info(f"Saved transcription for video '{video_title}' to S3")
                
            # Create default collection
            collection_id = await create_default_collection(
                material_type="transcriptions",
                material_id=job_id,
                material_metadata=metadata,
                current_user=current_user
            )
                
            return {
                "transcript": transcript_text,
                "metadata": metadata,
                "job_id": job_id,
                "collection_id": collection_id,
                "video_title": video_title
            }
        else:
            raise HTTPException(
                status_code=500,
                detail="Transcription job failed"
            )
            
    except Exception as e:
        logging.error(f"Error transcribing YouTube video: {str(e)}")
        # Clean up any temporary files if they exist
        if os.path.exists(f"{temp_audio_path[:-4]}.mp3"):
            os.remove(f"{temp_audio_path[:-4]}.mp3")
        if 's3_audio_key' in locals():
            s3_client.delete_object(Bucket=S3_BUCKET_NAME, Key=s3_audio_key)
        raise HTTPException(
            status_code=500,
            detail=f"Failed to transcribe video: {str(e)}"
        )

def extract_video_id(url: str) -> Optional[str]:
    """Extract YouTube video ID from various URL formats."""
    patterns = [
        r'(?:v=|\/)([0-9A-Za-z_-]{11}).*',  # Standard and shortened URLs
        r'(?:embed\/)([0-9A-Za-z_-]{11})',   # Embed URLs
        r'(?:youtu\.be\/)([0-9A-Za-z_-]{11})'  # Shortened URLs
    ]
    
    for pattern in patterns:
        match = re.search(pattern, url)
        if match:
            return match.group(1)
    return None

@app.get("/user-transcriptions")
async def get_user_transcriptions(current_user: str = Depends(get_current_user)):
    try:
        # List objects in the user's transcriptions metadata folder
        response = s3_client.list_objects_v2(
            Bucket=S3_BUCKET_NAME,
            Prefix=f"transcriptions/{current_user}/metadata/"
        )
        
        transcriptions = []
        for obj in response.get('Contents', []):
            metadata = json.loads(s3_client.get_object(Bucket=S3_BUCKET_NAME, Key=obj['Key'])['Body'].read())
            transcriptions.append({
                "job_id": metadata['job_id'],
                "video_title": metadata['video_title'],
                "video_url": metadata['video_url'],
                "transcription_date": metadata['transcription_date']
            })
        
        return transcriptions
    except Exception as e:
        logging.error(f"Failed to retrieve user transcriptions from S3: {str(e)}")
        raise HTTPException(status_code=500, detail="Failed to retrieve transcriptions")

@app.post("/transcribe-lecture")
async def transcribe_lecture(
    audio: UploadFile = File(...),
    title: str = Form(...),  # Make title a required field
    current_user: str = Depends(get_current_user)
):
    try:
        # Define supported formats
        SUPPORTED_FORMATS = {
            'audio': ['.mp3', '.wav', '.m4a'],
            'video': ['.mp4']
        }
        
        try:
            # Check file extension
            file_ext = os.path.splitext(audio.filename)[1].lower()
            all_supported_formats = SUPPORTED_FORMATS['audio'] + SUPPORTED_FORMATS['video']
            
            if file_ext not in all_supported_formats:
                raise HTTPException(
                    status_code=400,
                    detail=f"Unsupported file format. Supported formats: {', '.join(all_supported_formats)}"
                )

            # Generate unique IDs for the files
            job_id = str(uuid.uuid4())
            temp_input_path = f"/tmp/{job_id}_input{file_ext}"
            temp_output_path = f"/tmp/{job_id}.mp3"
            s3_audio_key = f"temp-audio/{job_id}.mp3"
            
            # Save the uploaded file temporarily
            with open(temp_input_path, "wb") as buffer:
                buffer.write(await audio.read())
            
            # Convert to mp3 if needed
            if file_ext != '.mp3':
                try:
                    subprocess.run([
                        'ffmpeg',
                        '-i', temp_input_path,
                        '-vn',  # Disable video if present
                        '-acodec', 'libmp3lame',
                        '-ab', '192k',
                        '-ar', '44100',
                        '-y',  # Overwrite output file if exists
                        temp_output_path
                    ], check=True, capture_output=True)
                    
                    # Remove input file after conversion
                    os.remove(temp_input_path)
                except subprocess.CalledProcessError as e:
                    logging.error(f"FFmpeg conversion error: {e.stderr.decode()}")
                    raise HTTPException(
                        status_code=500,
                        detail="Failed to convert audio format"
                    )
            else:
                # If it's already MP3, just rename
                os.rename(temp_input_path, temp_output_path)
            
            # Upload to S3
            s3_client.upload_file(temp_output_path, S3_BUCKET_NAME, s3_audio_key)
            s3_uri = f"s3://{S3_BUCKET_NAME}/{s3_audio_key}"
            
            # Start transcription job
            transcribe.start_transcription_job(
                TranscriptionJobName=job_id,
                Media={'MediaFileUri': s3_uri},
                MediaFormat='mp3',
                LanguageCode='en-US'
            )
            
            # Wait for transcription to complete
            while True:
                status = transcribe.get_transcription_job(TranscriptionJobName=job_id)
                if status['TranscriptionJob']['TranscriptionJobStatus'] in ['COMPLETED', 'FAILED']:
                    break
                await asyncio.sleep(5)
            
            # Clean up temporary files
            if os.path.exists(temp_output_path):
                os.remove(temp_output_path)
            s3_client.delete_object(Bucket=S3_BUCKET_NAME, Key=s3_audio_key)
            
            if status['TranscriptionJob']['TranscriptionJobStatus'] == 'COMPLETED':
                transcript_uri = status['TranscriptionJob']['Transcript']['TranscriptFileUri']
                
                # Get the transcript JSON
                async with httpx.AsyncClient() as client:
                    response = await client.get(transcript_uri)
                    transcript_data = response.json()
                
                # Create metadata and transcript objects for S3
                metadata = {
                    "file_name": audio.filename,
                    "title": title,  # Use provided title
                    "upload_date": datetime.now().isoformat(),
                    "file_type": "uploaded_audio" if file_ext in SUPPORTED_FORMATS['audio'] else "uploaded_video",
                    "job_id": job_id,
                    "original_filename": title  # Use title for collection naming
                }
                
                transcript_text = transcript_data['results']['transcripts'][0]['transcript']
                
                # Save metadata to S3
                metadata_key = f"transcriptions/{current_user}/metadata/{job_id}.json"
                s3_client.put_object(
                    Bucket=S3_BUCKET_NAME,
                    Key=metadata_key,
                    Body=json.dumps(metadata),
                    ContentType='application/json'
                )
                
                # Save transcript to S3
                transcript_key = f"transcriptions/{current_user}/content/{job_id}.txt"
                s3_client.put_object(
                    Bucket=S3_BUCKET_NAME,
                    Key=transcript_key,
                    Body=transcript_text,
                    ContentType='text/plain'
                )
                
                logging.info(f"Saved transcription for file '{audio.filename}' to S3")
                    
                # Create default collection
                collection_id = await create_default_collection(
                    material_type="transcriptions",
                    material_id=job_id,
                    material_metadata=metadata,
                    current_user=current_user
                )
                    
                return {
                    "transcript": transcript_text,
                    "metadata": metadata,
                    "job_id": job_id,
                    "collection_id": collection_id
                }
            else:
                raise HTTPException(
                    status_code=500,
                    detail="Transcription job failed"
                )
                
        except Exception as e:
            logging.error(f"Error transcribing audio file: {str(e)}")
            # Clean up any temporary files
            for path in [temp_input_path, temp_output_path]:
                if 'path' in locals() and os.path.exists(path):
                    os.remove(path)
            if 's3_audio_key' in locals():
                s3_client.delete_object(Bucket=S3_BUCKET_NAME, Key=s3_audio_key)
            raise HTTPException(
                status_code=500,
                detail=f"Failed to transcribe audio: {str(e)}"
            )

    except Exception as e:
        logging.error(f"Error transcribing lecture: {str(e)}")
        # Clean up any temporary files
        for path in [temp_input_path, temp_output_path]:
            if 'path' in locals() and os.path.exists(path):
                os.remove(path)
        if 's3_audio_key' in locals():
            s3_client.delete_object(Bucket=S3_BUCKET_NAME, Key=s3_audio_key)
        raise HTTPException(
            status_code=500,
            detail=f"Failed to transcribe lecture: {str(e)}"
        )

@app.post("/process-presentation")
async def process_presentation(
    presentation: UploadFile = File(...),
    current_user: str = Depends(get_current_user)
):
    if not presentation.filename.endswith('.pptx'):
        raise HTTPException(
            status_code=400,
            detail="Only .pptx files are supported"
        )

    try:
        # Generate unique ID for this presentation
        presentation_id = str(uuid.uuid4())
        temp_path = f"/tmp/{presentation_id}.pptx"

        # Save the uploaded file temporarily
        with open(temp_path, "wb") as buffer:
            buffer.write(await presentation.read())

        # Parse the presentation
        prs = Presentation(temp_path)
        
        # Extract content and structure
        slides_content = []
        for slide_number, slide in enumerate(prs.slides, 1):
            slide_content = {
                "number": slide_number,
                "title": "",
                "content": [],
                "notes": "",
            }

            # Get slide title if it exists
            if slide.shapes.title:
                slide_content["title"] = slide.shapes.title.text

            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content["content"].append(shape.text.strip())

            # Get speaker notes if they exist
            if slide.notes_slide and slide.notes_slide.notes_text_frame:
                slide_content["notes"] = slide.notes_slide.notes_text_frame.text

            slides_content.append(slide_content)

        # Create metadata
        metadata = {
            "presentation_id": presentation_id,
            "original_filename": presentation.filename,
            "upload_date": datetime.now().isoformat(),
            "total_slides": len(slides_content),
            "has_speaker_notes": any(slide["notes"] for slide in slides_content),
            "slide_titles": [slide["title"] for slide in slides_content if slide["title"]]
        }

        # Save metadata to S3
        metadata_key = f"presentations/{current_user}/metadata/{presentation_id}.json"
        s3_client.put_object(
            Bucket=S3_BUCKET_NAME,
            Key=metadata_key,
            Body=json.dumps(metadata),
            ContentType='application/json'
        )

        # Save each slide's content separately
        for slide in slides_content:
            # Combine all text content
            full_content = {
                "title": slide["title"],
                "content": slide["content"],
                "notes": slide["notes"]
            }
            
            slide_key = f"presentations/{current_user}/content/{presentation_id}/slide_{slide['number']}.json"
            s3_client.put_object(
                Bucket=S3_BUCKET_NAME,
                Key=slide_key,
                Body=json.dumps(full_content),
                ContentType='application/json'
            )

        # After successful processing, create a default collection
        collection_id = await create_default_collection(
            material_type="presentations",
            material_id=presentation_id,
            material_metadata=metadata,
            current_user=current_user
        )

        # Clean up temporary file
        os.remove(temp_path)

        logging.info(f"Successfully processed presentation: {presentation.filename}")
        
        return {
            "message": "Presentation processed successfully",
            "presentation_id": presentation_id,
            "metadata": metadata,
            "collection_id": collection_id  # Return the collection ID
        }

    except Exception as e:
        logging.error(f"Error processing presentation: {str(e)}")
        if 'temp_path' in locals() and os.path.exists(temp_path):
            os.remove(temp_path)
        raise HTTPException(
            status_code=500,
            detail=f"Failed to process presentation: {str(e)}"
        )

@app.get("/list-presentations/{user_id}")
async def list_presentations(
    user_id: str,
    current_user: str = Depends(get_current_user)
):
    if user_id != current_user:
        raise HTTPException(
            status_code=403,
            detail="You don't have permission to access this data"
        )
    
    try:
        # List all presentation metadata files for the user
        response = s3_client.list_objects_v2(
            Bucket=S3_BUCKET_NAME,
            Prefix=f"presentations/{user_id}/metadata/"
        )
        
        presentations = []
        for obj in response.get('Contents', []):
            metadata = json.loads(
                s3_client.get_object(
                    Bucket=S3_BUCKET_NAME,
                    Key=obj['Key']
                )['Body'].read()
            )
            presentations.append(metadata)
        
        return presentations
    
    except Exception as e:
        logging.error(f"Error listing presentations: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail="Failed to list presentations"
        )

@app.get("/get-presentation/{user_id}/{presentation_id}")
async def get_presentation(
    user_id: str,
    presentation_id: str,
    current_user: str = Depends(get_current_user),
    format: str = Query("separate", enum=["separate", "consolidated"])  # New query parameter
):
    if user_id != current_user:
        raise HTTPException(
            status_code=403,
            detail="You don't have permission to access this data"
        )
    
    try:
        # Get metadata
        metadata_key = f"presentations/{user_id}/metadata/{presentation_id}.json"
        metadata = json.loads(
            s3_client.get_object(
                Bucket=S3_BUCKET_NAME,
                Key=metadata_key
            )['Body'].read()
        )
        
        # Get all slides
        slides = []
        for slide_num in range(1, metadata['total_slides'] + 1):
            slide_key = f"presentations/{user_id}/content/{presentation_id}/slide_{slide_num}.json"
            slide_content = json.loads(
                s3_client.get_object(
                    Bucket=S3_BUCKET_NAME,
                    Key=slide_key
                )['Body'].read()
            )
            slides.append(slide_content)
        
        if format == "consolidated":
            # Create consolidated format
            consolidated_content = {
                "metadata": metadata,
                "content": {}
            }
            
            for slide_num, slide in enumerate(slides, 1):
                consolidated_content["content"][f"slide_{slide_num}"] = {
                    "title": slide["title"],
                    "content": slide["content"],
                    "notes": slide["notes"]
                }
            
            return consolidated_content
        else:
            # Return original separate format
            return {
                "metadata": metadata,
                "slides": slides
            }
        
    except Exception as e:
        logging.error(f"Error retrieving presentation: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail="Failed to retrieve presentation"
        )

@app.post("/collections")
async def create_collection(
    name: str = Body(...),
    materials: Dict = Body(...),  # Will contain textbook_sections, transcriptions, presentations
    current_user: str = Depends(get_current_user)
):
    try:
        collection_id = str(uuid.uuid4())
        collection = {
            "collection_id": collection_id,
            "name": name,
            "created_date": datetime.now().isoformat(),
            "user_id": current_user,
            "materials": materials
        }
        
        # Save to S3
        collection_key = f"collections/{current_user}/{collection_id}.json"
        s3_client.put_object(
            Bucket=S3_BUCKET_NAME,
            Key=collection_key,
            Body=json.dumps(collection),
            ContentType='application/json'
        )
        
        return collection

    except Exception as e:
        logging.error(f"Error creating collection: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail="Failed to create collection"
        )

@app.get("/collections")
async def list_collections(current_user: str = Depends(get_current_user)):
    try:
        # List all collections for the user
        response = s3_client.list_objects_v2(
            Bucket=S3_BUCKET_NAME,
            Prefix=f"collections/{current_user}/"
        )
        
        collections = []
        for obj in response.get('Contents', []):
            collection = json.loads(
                s3_client.get_object(
                    Bucket=S3_BUCKET_NAME,
                    Key=obj['Key']
                )['Body'].read()
            )
            
            # Count total materials in collection
            total_materials = sum(
                len(materials) for materials in collection['materials'].values()
            )
            
            # Only include collections with more than one material
            if total_materials > 1:
                collections.append(collection)
        
        return collections

    except Exception as e:
        logging.error(f"Error listing collections: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail="Failed to list collections"
        )

@app.put("/collections/{collection_id}/materials")
async def update_collection_materials(
    collection_id: str,
    materials: Dict = Body(...),
    current_user: str = Depends(get_current_user)
):
    try:
        # Get existing collection
        collection_key = f"collections/{current_user}/{collection_id}.json"
        collection = json.loads(
            s3_client.get_object(
                Bucket=S3_BUCKET_NAME,
                Key=collection_key
            )['Body'].read()
        )
        
        # Verify ownership
        if collection['user_id'] != current_user:
            raise HTTPException(status_code=403, detail="Not authorized to modify this collection")
        
        # Update materials
        collection['materials'] = materials
        
        # Save back to S3
        s3_client.put_object(
            Bucket=S3_BUCKET_NAME,
            Key=collection_key,
            Body=json.dumps(collection),
            ContentType='application/json'
        )
        
        return collection

    except ClientError as e:
        if e.response['Error']['Code'] == 'NoSuchKey':
            raise HTTPException(status_code=404, detail="Collection not found")
        raise HTTPException(status_code=500, detail="S3 error")
    except Exception as e:
        logging.error(f"Error updating collection: {str(e)}")
        raise HTTPException(status_code=500, detail="Failed to update collection")

@app.get("/collections/{collection_id}")
async def get_collection(
    collection_id: str,
    current_user: str = Depends(get_current_user)
):
    try:
        collection_key = f"collections/{current_user}/{collection_id}.json"
        collection = json.loads(
            s3_client.get_object(
                Bucket=S3_BUCKET_NAME,
                Key=collection_key
            )['Body'].read()
        )
        
        if collection['user_id'] != current_user:
            raise HTTPException(status_code=403, detail="Not authorized to access this collection")
            
        return collection

    except ClientError as e:
        if e.response['Error']['Code'] == 'NoSuchKey':
            raise HTTPException(status_code=404, detail="Collection not found")
        raise HTTPException(status_code=500, detail="S3 error")
    except Exception as e:
        logging.error(f"Error retrieving collection: {str(e)}")
        raise HTTPException(status_code=500, detail="Failed to retrieve collection")

@app.post("/process-notes")
async def process_notes(
    notes: UploadFile = File(...),
    current_user: str = Depends(get_current_user)
):
    SUPPORTED_FORMATS = ['.pdf', '.jpg', '.jpeg', '.png']
    
    try:
        # Check file format
        file_ext = os.path.splitext(notes.filename)[1].lower()
        if file_ext not in SUPPORTED_FORMATS:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file format. Supported formats: {', '.join(SUPPORTED_FORMATS)}"
            )

        # Generate unique ID
        notes_id = str(uuid.uuid4())
        temp_path = f"/tmp/{notes_id}{file_ext}"
        
        # Save file temporarily
        with open(temp_path, "wb") as buffer:
            buffer.write(await notes.read())
        
        # Upload to S3
        s3_key = f"notes/{current_user}/original/{notes_id}{file_ext}"
        s3_client.upload_file(temp_path, S3_BUCKET_NAME, s3_key)
        
        # Start Textract job
        if file_ext == '.pdf':
            textract_response = textract_client.start_document_analysis(
                DocumentLocation={'S3Object': {'Bucket': S3_BUCKET_NAME, 'Name': s3_key}},
                FeatureTypes=['TABLES', 'FORMS']
            )
            job_id = textract_response['JobId']
            
            # Wait for Textract job to complete
            while True:
                response = textract_client.get_document_analysis(JobId=job_id)
                if response['JobStatus'] in ['SUCCEEDED', 'FAILED']:
                    break
                await asyncio.sleep(5)
                
            if response['JobStatus'] == 'FAILED':
                raise Exception("Textract processing failed")
                
            textract_results = response
        else:
            # For images, use synchronous Textract
            with open(temp_path, 'rb') as image:
                textract_results = textract_client.detect_document_text(
                    Document={'Bytes': image.read()}
                )

        # Process with Rekognition for diagrams
        with open(temp_path, 'rb') as image:
            rekognition_results = rekognition.detect_labels(
                Image={'Bytes': image.read()},
                MaxLabels=10,
                MinConfidence=70
            )

        # Combine results
        processed_content = {
            "text_content": [],
            "diagrams": [],
            "tables": []
        }

        # Process Textract results
        if file_ext == '.pdf':
            for block in textract_results['Blocks']:
                if block['BlockType'] == 'LINE':
                    processed_content["text_content"].append({
                        "text": block['Text'],
                        "confidence": block['Confidence'],
                        "position": block['Geometry']
                    })
                elif block['BlockType'] == 'TABLE':
                    processed_content["tables"].append({
                        "position": block['Geometry'],
                        "confidence": block['Confidence']
                    })
        else:
            for item in textract_results['Blocks']:
                if item['BlockType'] == 'LINE':
                    processed_content["text_content"].append({
                        "text": item['Text'],
                        "confidence": item['Confidence'],
                        "position": item['Geometry']
                    })

        # Process Rekognition results
        for label in rekognition_results['Labels']:
            if label['Name'].lower() in ['diagram', 'chart', 'graph', 'drawing']:
                processed_content["diagrams"].append({
                    "type": label['Name'],
                    "confidence": label['Confidence'],
                    "position": label['Instances'][0]['BoundingBox'] if label['Instances'] else None
                })

        # Create metadata
        metadata = {
            "notes_id": notes_id,
            "original_filename": notes.filename,
            "upload_date": datetime.now().isoformat(),
            "file_type": file_ext,
            "processing_status": "completed",
            "has_diagrams": len(processed_content["diagrams"]) > 0,
            "has_tables": len(processed_content["tables"]) > 0
        }

        # Save metadata to S3
        metadata_key = f"notes/{current_user}/metadata/{notes_id}.json"
        s3_client.put_object(
            Bucket=S3_BUCKET_NAME,
            Key=metadata_key,
            Body=json.dumps(metadata),
            ContentType='application/json'
        )

        # Save processed content to S3
        content_key = f"notes/{current_user}/processed/{notes_id}.json"
        s3_client.put_object(
            Bucket=S3_BUCKET_NAME,
            Key=content_key,
            Body=json.dumps(processed_content),
            ContentType='application/json'
        )

        # Clean up temporary file
        os.remove(temp_path)

        # Create default collection
        collection_id = await create_default_collection(
            material_type="notes",
            material_id=notes_id,
            material_metadata=metadata,
            current_user=current_user
        )

        return {
            "message": "Notes processed successfully",
            "notes_id": notes_id,
            "metadata": metadata,
            "content": processed_content,
            "collection_id": collection_id
        }

    except Exception as e:
        logging.error(f"Error processing notes: {str(e)}")
        if 'temp_path' in locals() and os.path.exists(temp_path):
            os.remove(temp_path)
        raise HTTPException(
            status_code=500,
            detail=f"Failed to process notes: {str(e)}"
        )

@app.get("/notes/{notes_id}")
async def get_notes(
    notes_id: str,
    current_user: str = Depends(get_current_user)
):
    try:
        # Get metadata
        metadata_key = f"notes/{current_user}/metadata/{notes_id}.json"
        metadata = json.loads(
            s3_client.get_object(
                Bucket=S3_BUCKET_NAME,
                Key=metadata_key
            )['Body'].read()
        )

        # Get processed content
        content_key = f"notes/{current_user}/processed/{notes_id}.json"
        content = json.loads(
            s3_client.get_object(
                Bucket=S3_BUCKET_NAME,
                Key=content_key
            )['Body'].read()
        )

        return {
            "metadata": metadata,
            "content": content
        }

    except ClientError as e:
        if e.response['Error']['Code'] == 'NoSuchKey':
            raise HTTPException(status_code=404, detail="Notes not found")
        raise HTTPException(status_code=500, detail="S3 error")
    except Exception as e:
        logging.error(f"Error retrieving notes: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail="Failed to retrieve notes"
        )

async def get_collection_content(collection_id: str, current_user: str):
    try:
        # Get collection metadata
        collection = json.loads(
            s3_client.get_object(
                Bucket=S3_BUCKET_NAME,
                Key=f"collections/{current_user}/{collection_id}.json"
            )['Body'].read()
        )
        
        content = {
            "textbook_content": [],
            "transcriptions": [],
            "presentations": [],
            "notes": []
        }
        
        # Gather textbook sections
        for section in collection['materials'].get('textbook_sections', []):
            section_content = json.loads(
                s3_client.get_object(
                    Bucket=S3_BUCKET_NAME,
                    Key=f"textbooks/{current_user}/content/{section['book_id']}/section_{section['section_id']}.json"
                )['Body'].read()
            )
            content["textbook_content"].append(section_content)
        
        # Gather transcriptions
        for trans in collection['materials'].get('transcriptions', []):
            trans_content = json.loads(
                s3_client.get_object(
                    Bucket=S3_BUCKET_NAME,
                    Key=f"transcriptions/{current_user}/content/{trans['transcription_id']}.txt"
                )['Body'].read()
            )
            content["transcriptions"].append(trans_content)
        
        # Gather presentations
        for pres in collection['materials'].get('presentations', []):
            pres_content = json.loads(
                s3_client.get_object(
                    Bucket=S3_BUCKET_NAME,
                    Key=f"presentations/{current_user}/content/{pres['presentation_id']}"
                )['Body'].read()
            )
            content["presentations"].append(pres_content)
            
        # Gather notes
        for note in collection['materials'].get('notes', []):
            note_content = json.loads(
                s3_client.get_object(
                    Bucket=S3_BUCKET_NAME,
                    Key=f"notes/{current_user}/processed/{note['notes_id']}.json"
                )['Body'].read()
            )
            content["notes"].append(note_content)
            
        return content
        
    except Exception as e:
        logging.error(f"Error gathering collection content: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail="Failed to gather collection content"
        )

def generate_dynamic_prompt(content: Dict, task: str) -> str:
    # Base prompts for each task
    base_prompts = {
        "summary": """Generate a comprehensive summary that synthesizes information from multiple learning materials. 
                     Focus on key concepts, relationships between ideas, and important takeaways.""",
        
        "game": """Create an interactive learning game concept that tests understanding of the material. 
                  Include specific questions and scenarios based on the content.""",
        
        "diagrams": """Generate clear, informative diagrams that visualize key concepts and relationships 
                      from the learning materials."""
    }
    
    prompt_parts = [base_prompts[task]]
    
    # Add content-specific instructions
    if content["textbook_content"]:
        prompt_parts.append(
            "From the textbook sections, incorporate key definitions, concepts, and theoretical frameworks."
        )
    
    if content["transcriptions"]:
        prompt_parts.append(
            "From the lecture transcriptions, include practical examples, explanations, and real-world applications."
        )
    
    if content["presentations"]:
        prompt_parts.append(
            "From the presentation slides, use the main points, visual concepts, and structured progression of ideas."
        )
    
    if content["notes"]:
        prompt_parts.append(
            "From the handwritten notes, include additional insights, annotations, and supplementary examples."
        )
    
    # Add task-specific modifications
    if task == "summary":
        prompt_parts.append(
            "Create a coherent narrative that flows naturally between different source materials."
        )
    elif task == "game":
        prompt_parts.append(
            "Design interactions that test understanding across all available materials."
        )
    elif task == "diagrams":
        prompt_parts.append(
            "Create visualizations that show relationships between concepts from different sources."
        )
    
    # Combine all parts
    full_prompt = "\n\n".join(prompt_parts)
    
    # Add the actual content
    content_text = format_content_for_prompt(content)
    
    return f"{full_prompt}\n\nContent to work with:\n\n{content_text}"

def format_content_for_prompt(content: Dict) -> str:
    formatted_parts = []
    
    if content["textbook_content"]:
        formatted_parts.append("TEXTBOOK SECTIONS:")
        for section in content["textbook_content"]:
            formatted_parts.append(f"- {section['text']}")
    
    if content["transcriptions"]:
        formatted_parts.append("\nLECTURE TRANSCRIPTIONS:")
        for trans in content["transcriptions"]:
            formatted_parts.append(f"- {trans}")
    
    if content["presentations"]:
        formatted_parts.append("\nPRESENTATIONS:")
        for pres in content["presentations"]:
            slides = [f"Slide {i+1}: {slide['content']}" 
                     for i, slide in enumerate(pres['slides'])]
            formatted_parts.extend(slides)
    
    if content["notes"]:
        formatted_parts.append("\nHANDWRITTEN NOTES:")
        for note in content["notes"]:
            formatted_parts.append(f"- {note['text_content']}")
            if note['diagrams']:
                formatted_parts.append("  (Includes diagrams)")
    
    return "\n".join(formatted_parts)

# Helper function to create a default collection
async def create_default_collection(
    material_type: str,
    material_id: str,
    material_metadata: Dict,
    current_user: str
) -> str:
    try:
        collection_id = str(uuid.uuid4())
        collection = {
            "collection_id": collection_id,
            "name": f"{material_metadata.get('original_filename', 'Untitled')} Collection",
            "created_date": datetime.now().isoformat(),
            "user_id": current_user,
            "materials": {
                "textbook_sections": [],
                "transcriptions": [],
                "presentations": [],
                "notes": []
            }
        }
        
        # Add the material to its appropriate type in the collection
        collection["materials"][material_type].append({
            f"{material_type[:-1]}_id": material_id,  # Remove 's' from type for ID
            "added_date": datetime.now().isoformat()
        })
        
        # Save to S3
        collection_key = f"collections/{current_user}/{collection_id}.json"
        s3_client.put_object(
            Bucket=S3_BUCKET_NAME,
            Key=collection_key,
            Body=json.dumps(collection),
            ContentType='application/json'
        )
        
        return collection_id
        
    except Exception as e:
        logging.error(f"Error creating default collection: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail="Failed to create default collection"
        )
